import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Widescreen 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Theme Colors
    BG_COLOR = RGBColor(15, 23, 42)       # Slate 900
    CARD_BG_COLOR = RGBColor(30, 41, 59)  # Slate 800
    TEXT_MAIN = RGBColor(255, 255, 255)   # White
    TEXT_MUTED = RGBColor(148, 163, 184)  # Slate 400
    ACCENT_INDIGO = RGBColor(99, 102, 241)# Indigo 500
    ACCENT_TEAL = RGBColor(20, 184, 166)  # Teal 500

    # Fonts
    FONT_TITLE = "Trebuchet MS"
    FONT_BODY = "Arial"

    def apply_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_header(slide, title_text, category="MOCKAI PRESENTATION"):
        # Category Tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.73), Inches(0.4))
        cat_tf = cat_box.text_frame
        cat_tf.word_wrap = True
        cat_tf.margin_left = cat_tf.margin_right = cat_tf.margin_top = cat_tf.margin_bottom = 0
        cat_p = cat_tf.paragraphs[0]
        cat_p.text = category.upper()
        cat_p.font.name = FONT_TITLE
        cat_p.font.size = Pt(10)
        cat_p.font.bold = True
        cat_p.font.color.rgb = ACCENT_TEAL

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.73), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.margin_left = title_tf.margin_right = title_tf.margin_top = title_tf.margin_bottom = 0
        title_p = title_tf.paragraphs[0]
        title_p.text = title_text
        title_p.font.name = FONT_TITLE
        title_p.font.size = Pt(32)
        title_p.font.bold = True
        title_p.font.color.rgb = TEXT_MAIN

        # Decorative line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(2.0), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_INDIGO
        line.line.fill.background()

    def add_bullet_points(slide, left, top, width, height, points):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        for i, pt in enumerate(points):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.space_after = Pt(12)
            p.level = 0
            
            # Format title in bold, followed by details
            if "::" in pt:
                bold_part, regular_part = pt.split("::", 1)
                run1 = p.add_run()
                run1.text = "•  " + bold_part + ": "
                run1.font.name = FONT_BODY
                run1.font.size = Pt(15)
                run1.font.bold = True
                run1.font.color.rgb = ACCENT_TEAL
                
                run2 = p.add_run()
                run2.text = regular_part
                run2.font.name = FONT_BODY
                run2.font.size = Pt(15)
                run2.font.color.rgb = TEXT_MAIN
            else:
                run = p.add_run()
                run.text = "•  " + pt
                run.font.name = FONT_BODY
                run.font.size = Pt(15)
                run.font.color.rgb = TEXT_MAIN

    # --- SLIDE 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide1)

    # Accent decorative background card
    card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.5), Inches(11.33), Inches(4.5))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG_COLOR
    card.line.color.rgb = ACCENT_INDIGO
    card.line.width = Pt(1.5)

    # Title box inside card
    title_box = slide1.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.33), Inches(1.5))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "MockAI"
    p1.font.name = FONT_TITLE
    p1.font.size = Pt(64)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_MAIN

    p1_sub = tf1.add_paragraph()
    p1_sub.text = "AI-POWERED INTERVIEW PREPARATION ASSISTANT"
    p1_sub.font.name = FONT_TITLE
    p1_sub.font.size = Pt(18)
    p1_sub.font.bold = True
    p1_sub.font.color.rgb = ACCENT_TEAL
    p1_sub.space_before = Pt(8)

    # Details
    desc_box = slide1.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10.33), Inches(1.2))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    p_desc = tf_desc.paragraphs[0]
    p_desc.text = "A full-stack Generative AI application designed to conduct adaptive mock interviews, evaluate responses in real-time with granular feedback, map syllabus weak topics, and outline personalized learning roadmaps."
    p_desc.font.name = FONT_BODY
    p_desc.font.size = Pt(16)
    p_desc.font.color.rgb = TEXT_MUTED

    # --- SLIDE 2: The Problem ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide2)
    add_header(slide2, "The Interview Preparation Dilemma", "Market & Problem Context")
    
    # Left Box - Problem statement
    left_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = CARD_BG_COLOR
    left_box.line.color.rgb = ACCENT_INDIGO
    left_box.line.width = Pt(1.5)
    
    tb_l = slide2.shapes.add_textbox(Inches(1.1), Inches(2.3), Inches(4.9), Inches(3.9))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    pl1 = tf_l.paragraphs[0]
    pl1.text = "The Challenge"
    pl1.font.name = FONT_TITLE
    pl1.font.size = Pt(22)
    pl1.font.bold = True
    pl1.font.color.rgb = ACCENT_INDIGO
    pl1.space_after = Pt(14)
    
    problems = [
        "High Anxiety::Candidates face extreme nervousness during live interviews due to lack of standard mock options.",
        "Generic Material::Standard study platforms supply generic questions instead of role-targeted evaluations.",
        "Blind Spots::Candidates submit answers without understanding why they are incomplete or wrong.",
        "Vague Actions::Standard evaluations rank candidates but fail to outline a step-by-step roadmap to study."
    ]
    add_bullet_points(slide2, Inches(1.1), Inches(2.9), Inches(4.9), Inches(3.3), problems)

    # Right Box - Impact statement
    right_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.5))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = CARD_BG_COLOR
    right_box.line.color.rgb = ACCENT_TEAL
    right_box.line.width = Pt(1.5)
    
    tb_r = slide2.shapes.add_textbox(Inches(7.3), Inches(2.3), Inches(4.9), Inches(3.9))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    pr1 = tf_r.paragraphs[0]
    pr1.text = "The Cost of Poor Prep"
    pr1.font.name = FONT_TITLE
    pr1.font.size = Pt(22)
    pr1.font.bold = True
    pr1.font.color.rgb = ACCENT_TEAL
    pr1.space_after = Pt(14)
    
    impacts = [
        "Missed Career Openings::Candidates fail interviews not from a lack of technical knowledge, but due to poor response structuring.",
        "Wasted Practice Hours::Practicing without target alignment wastes candidate time.",
        "Delayed Feedback Loops::Traditional mock portals require manual, expensive evaluations with long turnaround times."
    ]
    add_bullet_points(slide2, Inches(7.3), Inches(2.9), Inches(4.9), Inches(3.3), impacts)

    # --- SLIDE 3: The Solution ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide3)
    add_header(slide3, "MockAI: The Intelligent Solution", "The Value Proposition")

    # We will lay out 3 vertical feature cards
    card_width = Inches(3.64)
    card_height = Inches(4.5)
    gap = Inches(0.4)
    start_left = Inches(0.8)
    
    solutions = [
        {
            "title": "Interactive Simulation",
            "accent": ACCENT_INDIGO,
            "points": [
                "Tailored Mock Sessions::Configures custom sessions for 6 core professional roles.",
                "Hybrid Question Styles::Tests candidates using descriptive, MCQs, and fill-in-the-blank question sets.",
                "Flexible Settings::Allows Easy, Medium, and Hard filters across technical and behavioral categories."
            ]
        },
        {
            "title": "Real-time AI Grading",
            "accent": ACCENT_TEAL,
            "points": [
                "Structured Evaluation::Generates detailed score profiles across technical, communication, and relevance metrics.",
                "Granular Feedback::Returns specific notes explaining gaps and recommending improvements.",
                "Instant Analytics::Eliminates waiting time with live score calculators."
            ]
        },
        {
            "title": "Syllabus Roadmaps",
            "accent": ACCENT_INDIGO,
            "points": [
                "Weak Topic Mapping::Identifies gaps and links candidate errors to specific syllabus modules.",
                "Interactive Checklist::Autogenerates a personalized, checkable study plan.",
                "Session Resumption::Allows logging back in to review logs or continue incomplete mocks."
            ]
        }
    ]

    for idx, sol in enumerate(solutions):
        c_left = start_left + idx * (card_width + gap)
        # Card Shape
        card_shape = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c_left, Inches(2.0), card_width, card_height)
        card_shape.fill.solid()
        card_shape.fill.fore_color.rgb = CARD_BG_COLOR
        card_shape.line.color.rgb = sol["accent"]
        card_shape.line.width = Pt(1.5)
        
        # Header inside Card
        t_box = slide3.shapes.add_textbox(c_left + Inches(0.2), Inches(2.2), card_width - Inches(0.4), Inches(0.6))
        t_tf = t_box.text_frame
        t_tf.word_wrap = True
        tp = t_tf.paragraphs[0]
        tp.text = sol["title"]
        tp.font.name = FONT_TITLE
        tp.font.size = Pt(20)
        tp.font.bold = True
        tp.font.color.rgb = sol["accent"]
        
        # Bullets
        add_bullet_points(slide3, c_left + Inches(0.2), Inches(2.9), card_width - Inches(0.4), Inches(3.4), sol["points"])

    # --- SLIDE 4: Architecture ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide4)
    add_header(slide4, "Full-Stack System Architecture", "System Design & Data Flow")

    # Left: Text outline of components
    add_bullet_points(slide4, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.3), [
        "React Frontend (Vite)::Modern SPA using React 19, TypeScript, and TailwindCSS v4. Provides a fully-styled dark-mode dashboard with real-time state tracking.",
        "FastAPI Backend::Python 3.13 REST API endpoints. Utilizes Pydantic schemas, dependency injection security layers, and asynchronous query routers.",
        "MongoDB (Local Instance)::Stores secure user records, candidate settings profiles, historical mock reports, and the central seeded question bank.",
        "Google Gemini API (`gemini-1.5-flash`)::Drives evaluation scoring and generates personalized improvement roadmap checklists via JSON Structured Output Mode."
    ])

    # Right: A visual block layout using Shapes representing Architecture flow
    # Frontend Box
    fe_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(2.2), Inches(4.5), Inches(0.9))
    fe_box.fill.solid()
    fe_box.fill.fore_color.rgb = CARD_BG_COLOR
    fe_box.line.color.rgb = ACCENT_TEAL
    fe_tf = fe_box.text_frame
    fe_tf.word_wrap = True
    fe_p = fe_tf.paragraphs[0]
    fe_p.alignment = PP_ALIGN.CENTER
    fe_p.text = "React Frontend\n(TypeScript + Tailwind v4)"
    fe_p.font.name = FONT_TITLE
    fe_p.font.size = Pt(14)
    fe_p.font.bold = True
    fe_p.font.color.rgb = TEXT_MAIN

    # Arrow 1
    arr1 = slide4.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(9.55), Inches(3.2), Inches(0.4), Inches(0.4))
    arr1.fill.solid()
    arr1.fill.fore_color.rgb = ACCENT_INDIGO
    arr1.line.fill.background()

    # Backend Box
    be_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(3.7), Inches(4.5), Inches(0.9))
    be_box.fill.solid()
    be_box.fill.fore_color.rgb = CARD_BG_COLOR
    be_box.line.color.rgb = ACCENT_INDIGO
    be_tf = be_box.text_frame
    be_tf.word_wrap = True
    be_p = be_tf.paragraphs[0]
    be_p.alignment = PP_ALIGN.CENTER
    be_p.text = "FastAPI Backend\n(REST API Core Services)"
    be_p.font.name = FONT_TITLE
    be_p.font.size = Pt(14)
    be_p.font.bold = True
    be_p.font.color.rgb = TEXT_MAIN

    # Three horizontal columns on bottom
    sub_width = Inches(1.3)
    sub_gap = Inches(0.3)
    sub_y = Inches(5.3)
    
    sub_systems = [
        {"name": "Auth\n(JWT/Bcrypt)", "color": ACCENT_TEAL},
        {"name": "Database\n(MongoDB)", "color": ACCENT_INDIGO},
        {"name": "GenAI\n(Gemini API)", "color": ACCENT_TEAL}
    ]
    
    # Arrow 2 split
    arr2 = slide4.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(9.55), Inches(4.7), Inches(0.4), Inches(0.4))
    arr2.fill.solid()
    arr2.fill.fore_color.rgb = ACCENT_INDIGO
    arr2.line.fill.background()

    for idx, sys_item in enumerate(sub_systems):
        x = Inches(7.5) + idx * (sub_width + sub_gap)
        sb = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, sub_y, sub_width, Inches(0.8))
        sb.fill.solid()
        sb.fill.fore_color.rgb = CARD_BG_COLOR
        sb.line.color.rgb = sys_item["color"]
        sb_tf = sb.text_frame
        sb_tf.word_wrap = True
        sb_p = sb_tf.paragraphs[0]
        sb_p.alignment = PP_ALIGN.CENTER
        sb_p.text = sys_item["name"]
        sb_p.font.name = FONT_BODY
        sb_p.font.size = Pt(11)
        sb_p.font.bold = True
        sb_p.font.color.rgb = TEXT_MAIN

    # --- SLIDE 5: Seeding & Questions ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide5)
    add_header(slide5, "Custom Seeding & Structured Bank", "Question Database Configuration")

    add_bullet_points(slide5, Inches(0.8), Inches(2.2), Inches(11.73), Inches(4.5), [
        "Pre-Configured Seed Script::Integrates a seed utility (`questions_bank.py`) that loads on first database boot, populating MongoDB instantly with clean datasets.",
        "540 Custom Mock Questions::Includes detailed datasets covering 6 developer roles, 3 interview styles, and 3 difficulty settings.",
        "Granular Matrix Coverage::Divides configurations into 10 unique sheets containing 10 questions per sheet, preventing redundant questions.",
        "Flexible Question Models::Accommodates multiple types:\n   - Multiple Choice Questions (MCQ) containing 4 option tags and correct answer schemas\n   - Fill-in-the-Blank (FIB) keys mapping targeted keywords\n   - Comprehensive Descriptive questions requesting detailed replies",
        "Fast Document Retrieval::Fetches active session questions directly from MongoDB matching specific setup wizard inputs, maintaining sub-millisecond retrieval speeds."
    ])

    # --- SLIDE 6: AI Evaluation Engine ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide6)
    add_header(slide6, "Structured AI Grading Engine", "Advanced Generative AI Service")

    # Left: Explanation of AI Evaluation
    add_bullet_points(slide6, Inches(0.8), Inches(2.2), Inches(6.0), Inches(4.5), [
        "Structured Schema Output::Commands the Gemini API using strict JSON schemas to output scores and feedback, parsing seamlessly into backend database structures.",
        "Evaluates 3 Core Metrics::\n   1. Technical Correctness (Accuracy of concepts, facts, syntax)\n   2. Communication Clarity (Sentence structure, articulation quality)\n   3. Answer Relevance (Direct response to constraints)",
        "Hybrid Local Parser Fallback::Features a local fallback engine that scores answers without an API key by matching keywords (FIB), validating option keys (MCQs), and analyzing word lengths (Descriptive)."
    ])

    # Right: Decorative Schema Preview Card
    schema_card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(2.2), Inches(5.0), Inches(4.3))
    schema_card.fill.solid()
    schema_card.fill.fore_color.rgb = CARD_BG_COLOR
    schema_card.line.color.rgb = ACCENT_TEAL
    schema_card.line.width = Pt(1.5)

    sc_tb = slide6.shapes.add_textbox(Inches(7.7), Inches(2.4), Inches(4.6), Inches(3.9))
    sc_tf = sc_tb.text_frame
    sc_tf.word_wrap = True
    scp = sc_tf.paragraphs[0]
    scp.text = "Gemini API Response Schema (Pydantic)"
    scp.font.name = FONT_TITLE
    scp.font.size = Pt(14)
    scp.font.bold = True
    scp.font.color.rgb = ACCENT_TEAL
    scp.space_after = Pt(10)

    schema_code = (
        "class AnswerEvaluation(BaseModel):\n"
        "    score_correctness: int  # 0-100\n"
        "    score_clarity: int      # 0-100\n"
        "    score_relevance: int    # 0-100\n"
        "    overall_score: int      # Average\n"
        "    feedback: str           # Detailed text\n"
        "    weak_topics: List[str]  # e.g., 'REST APIs'\n"
        "    model_answer: str       # Ideal response"
    )
    
    scp2 = sc_tf.add_paragraph()
    scp2.text = schema_code
    scp2.font.name = "Consolas"
    scp2.font.size = Pt(11)
    scp2.font.color.rgb = TEXT_MAIN
    scp2.space_before = Pt(4)

    # --- SLIDE 7: Live Interview Room ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide7)
    add_header(slide7, "Interactive Live Interview Room", "Interactive Frontend Experience")

    add_bullet_points(slide7, Inches(0.8), Inches(2.2), Inches(11.73), Inches(4.5), [
        "Tailored Response Widgets::Swaps textarea blocks for 4 premium clickable option cards during MCQs, renders compact input text fields for FIBs, and opens rich text boxes for descriptive topics.",
        "Interactive Word Suggestions::Displays dynamic helper counters recommending word targets based on difficulty (e.g. min 25 words for Technical Easy, min 50 for hard challenges).",
        "Session State Saver::Caches progress metrics at the database tier on each submission, allowing candidates to pause or resume sessions safely.",
        "Interactive Loader Framework::Displays loading skeletons during GenAI evaluation calls to indicate background scoring operations.",
        "10-Question Progress Indicators::Renders real-time trackers displaying active indexes, scores, and completion progress bars."
    ])

    # --- SLIDE 8: Detailed Performance Reports ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide8)
    add_header(slide8, "Granular Performance Scorecards", "Post-Interview Analytics & Roadmaps")

    # Left Box - SVG Score dials representation
    l_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5))
    l_box.fill.solid()
    l_box.fill.fore_color.rgb = CARD_BG_COLOR
    l_box.line.color.rgb = ACCENT_INDIGO
    l_box.line.width = Pt(1.5)
    
    tb_l8 = slide8.shapes.add_textbox(Inches(1.1), Inches(2.3), Inches(4.9), Inches(3.9))
    tf_l8 = tb_l8.text_frame
    tf_l8.word_wrap = True
    pl8 = tf_l8.paragraphs[0]
    pl8.text = "Analytics Dashboards"
    pl8.font.name = FONT_TITLE
    pl8.font.size = Pt(22)
    pl8.font.bold = True
    pl8.font.color.rgb = ACCENT_INDIGO
    pl8.space_after = Pt(14)
    
    analytics_points = [
        "Score Gauges::Circular SVG gauges trace percentage scores across correctness, clarity, and relevance.",
        "Topic Gap Spotter::Summarizes all weak syllabus areas flagged during the interview.",
        "Model Answer Compare::Displays user inputs alongside ideal model answers side-by-side."
    ]
    add_bullet_points(slide8, Inches(1.1), Inches(2.9), Inches(4.9), Inches(3.3), analytics_points)

    # Right Box - Improvement Roadmap
    r_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.5))
    r_box.fill.solid()
    r_box.fill.fore_color.rgb = CARD_BG_COLOR
    r_box.line.color.rgb = ACCENT_TEAL
    r_box.line.width = Pt(1.5)
    
    tb_r8 = slide8.shapes.add_textbox(Inches(7.3), Inches(2.3), Inches(4.9), Inches(3.9))
    tf_r8 = tb_r8.text_frame
    tf_r8.word_wrap = True
    pr8 = tf_r8.paragraphs[0]
    pr8.text = "Personalized Improvement Plan"
    pr8.font.name = FONT_TITLE
    pr8.font.size = Pt(22)
    pr8.font.bold = True
    pr8.font.color.rgb = ACCENT_TEAL
    pr8.space_after = Pt(14)
    
    roadmap_points = [
        "Custom Checklist Generator::Creates a customized, checkable action plan tailored to identified gaps.",
        "Structured Roadmaps::Aggregates questions failed, maps them to syllabus concepts, and outputs step-by-step reading roadmaps.",
        "Database History Tracker::Saves completed scorecard reports to the profile timeline for tracking performance curves."
    ]
    add_bullet_points(slide8, Inches(7.3), Inches(2.9), Inches(4.9), Inches(3.3), roadmap_points)

    # --- SLIDE 9: Endpoint Architecture ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide9)
    add_header(slide9, "Comprehensive Unified API Schema", "REST API System Specs")

    add_bullet_points(slide9, Inches(0.8), Inches(2.2), Inches(11.73), Inches(4.5), [
        "Authentication Controllers::Handles user management:\n   - `POST /api/auth/register` (creates new profiles)\n   - `POST /api/auth/login` (authenticates credentials and returns JWT bearer tokens)\n   - `GET /api/auth/me` (verifies session tokens)",
        "Profile Integrations::Manages settings:\n   - `GET /api/profile` (fetches educational tags, role metrics, and custom skills)\n   - `PUT /api/profile` (updates experience level settings and interactive tag selections)",
        "Interview Lifecycle Handlers::Controls mock progress:\n   - `POST /api/interview/start` (creates active sessions and returns database UUIDs)\n   - `GET /api/interview/question/{id}` (retrieves the active question index)\n   - `POST /api/interview/answer/{id}` (processes inputs and compiles intermediate scoring summaries)\n   - `POST /api/interview/complete/{id}` (aggregates weak areas and builds personalized roadmaps)",
        "Log Reports Query Engine::Manages historical data:\n   - `GET /api/history` (gathers completed session records for profile widgets)\n   - `GET /api/results/{id}` (retrieves full scorecard details)"
    ])

    # --- SLIDE 10: Future Roadmap & Conclusion ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide10)
    add_header(slide10, "Future Roadmap & Technical Extensions", "The Development Path Ahead")

    add_bullet_points(slide10, Inches(0.8), Inches(2.2), Inches(11.73), Inches(4.5), [
        "Voice-to-Text Integration::Incorporate browser-native Web Speech API tools to allow candidates to speak answers, grading speech-to-text transcriptions for grammar and tone.",
        "AI Resume Parser API::Create standard backend parsing modules (using PyMuPDF and LLMs) to upload PDF resumes and auto-populate candidate profile tags and skill selectors.",
        "Video Engagement Analytics::Implement client-side webcam analysis to track eye focus, speech pacing, and expression profiles, delivering non-verbal behavior dashboards.",
        "Collaborative Mock Rooms::Develop real-time peer review mock modules using WebSocket backends, enabling team dashboards for classrooms and recruiting panels.",
        "Enterprise Dashboard Core::Scale database schemas to support recruiter dashboards, allowing universities and coding bootcamps to track collective training graphs."
    ])

    # Save presentation
    output_path = "MockAI_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {os.path.abspath(output_path)}")

if __name__ == "__main__":
    create_presentation()
